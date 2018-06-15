from pptx import Presentation
from pptx.util import Inches, Cm


# This is the Standard Title and Content layout slide from the VMware dark 16x9 template theme
SLD_LAYOUT_TITLE_AND_CONTENT = 7
MSO_THEME_COLOR = 1
# Open Presentation template file
prs = Presentation('/Users/bohlea/Documents/GitHub/RoadmapPreso/vmc-on-aws-roadmap-template.pptx')


# Add a new slide to the presentation
slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
roadmap_service_slide = prs.slides.add_slide(slide_layout)
shapes = roadmap_service_slide.shapes
title = roadmap_service_slide.shapes.title
title.text = "Roadmap: Service"



rows = 7
cols = 4
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(0.8)

table = shapes.add_table(rows, cols, left, top, width, height).table


# table.fill.fore_color.theme_color == MSO_THEME_COLOR.ACCENT_1

# set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.0)

# write column headings
table.cell(0, 0).text = 'Key Themes'
table.cell(1, 0).text = 'AWS regions'
table.cell(2, 0).text = 'Consumption'
table.cell(3, 0).text = 'Support & Customer Success'
table.cell(4, 0).text = 'SLA & Maintenance'
table.cell(5, 0).text = 'Onboarding'
table.cell(6, 0).text = 'Compliance'



prs.save('/Users/bohlea/Documents/GitHub/RoadmapPreso/new-vmc-roadmap.pptx')